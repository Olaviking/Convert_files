from flask import Flask, request, jsonify
import os
import csv
from PyPDF2 import PdfReader
import openpyxl
from docx import Document
import requests
from bs4 import BeautifulSoup
from pptx import Presentation
import json
import xml.etree.ElementTree as ET
from youtube_transcript_api import YouTubeTranscriptApi
from youtube_transcript_api.formatters import JSONFormatter


app = Flask(__name__)
API_KEY = "placeholderAPIKey"

@app.route('/convert', methods=['POST'])
def convert():
    # 2. Check for an API key in the request headers.
    provided_api_key = request.headers.get('x-api-key')
    if not provided_api_key or provided_api_key != API_KEY:
        return jsonify(error="Invalid or missing API key"), 401
    
    if 'file' not in request.files and 'url' not in request.form:
        return jsonify(error="No file or URL provided"), 400
    
    if 'file' in request.files:
        file = request.files['file']
        if file.filename == '':
            return jsonify(error="No selected file"), 400

        filename = os.path.join("/tmp", file.filename)
        file.save(filename)
        text, error = extract_text_from_file(filename)
        os.remove(filename)
        if error:
            return jsonify(error=error), 500
        return jsonify(text=text)
    
    if 'url' in request.form:
        url = request.form['url']
        text, error = extract_text_from_url(url)
        if error:
            return jsonify(error=error), 500
        return jsonify(text=text)

def extract_text_from_file(filename):
    _, extension = os.path.splitext(filename)
    try:
        if extension == '.pdf':
            return extract_from_pdf(filename), None
        elif extension == '.docx':
            return extract_from_word(filename), None
        elif extension == '.xlsx':
            return extract_from_excel(filename), None
        elif extension == '.csv':
            return extract_from_csv(filename), None
        elif extension == '.txt':
            return extract_from_txt(filename), None
        elif extension == '.pptx':
            return extract_from_pptx(filename), None
        elif extension == '.html':
            return extract_from_html(filename), None
        elif extension == '.xml':
            return extract_from_xml(filename), None
        elif extension == '.json':
            return extract_from_json(filename), None
        else:
            return None, "Unsupported filetype"
    except Exception as e:
        return None, str(e)

def extract_from_pdf(filename):
    with open(filename, 'rb') as pdf_file:
        reader = PdfReader(pdf_file)
        text = " ".join(page.extract_text() for page in reader.pages)
    return text

def extract_from_word(filename):
    doc = Document(filename)
    return " ".join(paragraph.text for paragraph in doc.paragraphs)

def extract_from_excel(filename):
    wb = openpyxl.load_workbook(filename)
    all_text = ""
    for sheet in wb:
        for row in sheet.iter_rows():
            for cell in row:
                all_text += str(cell.value) + " "
    return all_text

def extract_from_csv(filename):
    all_text = ""
    with open(filename, 'r') as csv_file:
        reader = csv.reader(csv_file)
        for row in reader:
            all_text += " ".join(row) + " "
    return all_text

def extract_from_txt(filename):
    with open(filename, 'r') as txt_file:
        return txt_file.read()

def extract_from_pptx(filename):
    prs = Presentation(filename)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "
    return text

def extract_from_html(filename):
    with open(filename, 'r', encoding='utf-8') as html_file:
        soup = BeautifulSoup(html_file, 'html.parser')
        return " ".join(soup.stripped_strings)

def extract_from_xml(filename):
    tree = ET.parse(filename)
    root = tree.getroot()
    texts = [elem.text for elem in root.iter() if elem.text]
    return " ".join(texts)

def extract_from_json(filename):
    with open(filename, 'r', encoding='utf-8') as json_file:
        data = json.load(json_file)
        return json.dumps(data)

def extract_text_from_url(url):
    try:
        # Check if the URL is a YouTube video URL
        if "youtube.com/watch?v=" in url:
            video_id = url.split("v=")[1].split("&")[0]
            return extract_youtube_transcript(video_id)
        
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
        }
        response = requests.get(url, headers=headers)
        response.raise_for_status()
        soup = BeautifulSoup(response.content, 'html.parser')
        return " ".join(soup.stripped_strings), None
    except Exception as e:
        return None, str(e)


def extract_youtube_transcript(video_id):
    try:
        transcript = YouTubeTranscriptApi.get_transcript(video_id)
        # Extract the 'text' from each entry in the transcript and join them
        formatted_transcript = " ".join(entry['text'] for entry in transcript)
        return formatted_transcript, None
    except Exception as e:
        return None, str(e)



if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=8080)